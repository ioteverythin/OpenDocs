"""Export tools — convert Markdown drafts to Word, PDF, PPTX.

Prefers the OpenDocs pipeline when available, with built-in fallbacks.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from ..models.document_model import DraftDocument, ExportFormat

logger = logging.getLogger("docagent.tools.export")


class ExportTools:
    """Convert finalised Markdown drafts to output files."""

    def __init__(self, outputs_dir: Path, *, theme_name: str = "corporate") -> None:
        self._outputs_dir = outputs_dir
        self._outputs_dir.mkdir(parents=True, exist_ok=True)
        self._theme_name = theme_name

    # ------------------------------------------------------------------
    # Public
    # ------------------------------------------------------------------
    def export(
        self,
        draft: DraftDocument,
        fmt: ExportFormat,
        *,
        diagram_paths: dict[str, Any] | None = None,
    ) -> Path:
        """Export *draft* to the requested format and return the output path."""
        self._diagram_paths = diagram_paths or {}
        dispatch = {
            ExportFormat.WORD: self._to_word,
            ExportFormat.PDF: self._to_pdf,
            ExportFormat.PPTX: self._to_pptx,
        }
        handler = dispatch.get(fmt)
        if handler is None:
            raise ValueError(f"Unsupported export format: {fmt}")
        return handler(draft)

    def export_all(self, draft: DraftDocument, *, diagram_paths: dict[str, Any] | None = None) -> list[Path]:
        """Export a draft to all supported formats."""
        paths: list[Path] = []
        for fmt in ExportFormat:
            try:
                paths.append(self.export(draft, fmt, diagram_paths=diagram_paths))
            except Exception as exc:
                logger.warning("Export %s failed for %s: %s", fmt.value, draft.doc_type.value, exc)
        return paths

    # ------------------------------------------------------------------
    # Word (.docx)
    # ------------------------------------------------------------------
    def _to_word(self, draft: DraftDocument) -> Path:
        """Export to Word via OpenDocs pipeline, or fallback to python-docx."""
        out_path = self._outputs_dir / f"{draft.doc_type.value}.docx"

        # Try OpenDocs first
        if self._try_opendocs_word(draft, out_path):
            return out_path

        # Fallback: python-docx
        return self._fallback_word(draft, out_path)

    def _try_opendocs_word(self, draft: DraftDocument, out_path: Path) -> bool:
        """Attempt to use the OpenDocs Word generator."""
        try:
            from opendocs.core.parser import ReadmeParser
            from opendocs.generators.word_generator import WordGenerator
            from opendocs.generators.themes import get_theme

            parser = ReadmeParser()
            doc_model = parser.parse(
                draft.content,
                repo_name=draft.title,
                repo_url="",
            )
            theme = get_theme(self._theme_name)
            gen = WordGenerator(theme=theme)
            result = gen.generate(doc_model, out_path.parent)
            if result.success and result.output_path:
                # Rename to our target name
                result_path = Path(result.output_path)
                if result_path != out_path and result_path.exists():
                    result_path.rename(out_path)
                # Post-process: embed diagram PNGs replacing [Image: ...] placeholders
                self._embed_diagrams_in_docx(out_path)
                logger.info("Word exported via OpenDocs: %s", out_path)
                return True
        except Exception as exc:
            logger.debug("OpenDocs Word export failed, using fallback: %s", exc)
        return False

    def _embed_diagrams_in_docx(self, docx_path: Path) -> None:
        """Post-process a .docx — replace [Image: ...] placeholders with real PNGs."""
        diagram_map = self._diagram_paths
        if not diagram_map:
            return

        try:
            from docx import Document as DocxDocument
            from docx.shared import Inches
        except ImportError:
            return

        # Build a lookup: normalised alt text → Path
        label_to_path: dict[str, Path] = {}
        _label_map = {
            "architecture": "System Architecture",
            "data_flow": "Data Flow",
            "component": "Component Overview",
        }
        for dtype, dpath in diagram_map.items():
            if dpath is None:
                continue
            p = Path(str(dpath))
            if not p.exists():
                continue
            nice = _label_map.get(dtype, dtype)
            label_to_path[nice.lower()] = p
            # Also match by raw diagram type
            label_to_path[dtype.lower()] = p

        if not label_to_path:
            return

        placeholder_re = re.compile(r'^\[Image:\s*(.+?)\]$')

        doc = DocxDocument(str(docx_path))
        replaced = 0

        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            m = placeholder_re.match(text)
            if not m:
                continue
            alt = m.group(1).strip()
            img_path = label_to_path.get(alt.lower())
            if img_path is None:
                # Try partial match
                for key, path in label_to_path.items():
                    if key in alt.lower() or alt.lower() in key:
                        img_path = path
                        break
            if img_path is None or not img_path.exists():
                continue

            # Clear the placeholder text and embed the image
            try:
                para.clear()
                run = para.add_run()
                run.add_picture(str(img_path), width=Inches(5.5))
                replaced += 1
                logger.info("Embedded diagram in docx: %s", alt)
            except Exception as exc:
                # Restore placeholder if embedding fails
                para.clear()
                para.add_run(f"[Image: {alt}]").italic = True
                logger.warning("Failed to embed diagram %s: %s", alt, exc)

        if replaced:
            doc.save(str(docx_path))
            logger.info("Post-processed docx: embedded %d diagram(s)", replaced)

    def _fallback_word(self, draft: DraftDocument, out_path: Path) -> Path:
        """Simple python-docx fallback with diagram image embedding."""
        from docx import Document as DocxDocument
        from docx.shared import Inches

        doc = DocxDocument()

        # Regex to detect Markdown image references: ![alt](path)
        img_pattern = re.compile(r'^!\[([^\]]*)\]\(([^)]+)\)$')

        for line in draft.content.splitlines():
            stripped = line.strip()

            # Check for image reference
            img_match = img_pattern.match(stripped)
            if img_match:
                alt_text = img_match.group(1)
                img_path = img_match.group(2)
                try:
                    from pathlib import Path as P
                    p = P(img_path)
                    if p.exists() and p.suffix.lower() in ('.png', '.jpg', '.jpeg', '.gif'):
                        doc.add_picture(str(p), width=Inches(5.5))
                        if alt_text:
                            doc.add_paragraph(alt_text, style='Caption')
                        logger.info("Embedded diagram: %s", alt_text or img_path)
                        continue
                except Exception as exc:
                    logger.debug("Could not embed image %s: %s", img_path, exc)

            if stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            elif stripped.startswith("1. "):
                doc.add_paragraph(stripped[3:], style="List Number")
            elif stripped:
                doc.add_paragraph(stripped)

        doc.save(str(out_path))
        logger.info("Word exported (fallback): %s", out_path)
        return out_path

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------
    def _to_pdf(self, draft: DraftDocument) -> Path:
        out_path = self._outputs_dir / f"{draft.doc_type.value}.pdf"

        # Try OpenDocs
        if self._try_opendocs_pdf(draft, out_path):
            return out_path

        # Fallback: reportlab
        return self._fallback_pdf(draft, out_path)

    def _try_opendocs_pdf(self, draft: DraftDocument, out_path: Path) -> bool:
        try:
            from opendocs.core.parser import ReadmeParser
            from opendocs.generators.pdf_generator import PDFGenerator
            from opendocs.generators.themes import get_theme

            parser = ReadmeParser()
            doc_model = parser.parse(draft.content, repo_name=draft.title)
            theme = get_theme(self._theme_name)
            gen = PDFGenerator(theme=theme)
            result = gen.generate(doc_model, out_path.parent)
            if result.success and result.output_path:
                result_path = Path(result.output_path)
                if result_path != out_path and result_path.exists():
                    result_path.rename(out_path)
                logger.info("PDF exported via OpenDocs: %s", out_path)
                return True
        except Exception as exc:
            logger.debug("OpenDocs PDF failed, using fallback: %s", exc)
        return False

    def _fallback_pdf(self, draft: DraftDocument, out_path: Path) -> Path:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

        doc = SimpleDocTemplate(str(out_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story: list = []

        for line in draft.content.splitlines():
            stripped = line.strip()
            if stripped.startswith("# "):
                story.append(Paragraph(stripped[2:], styles["Title"]))
                story.append(Spacer(1, 12))
            elif stripped.startswith("## "):
                story.append(Paragraph(stripped[3:], styles["Heading2"]))
                story.append(Spacer(1, 8))
            elif stripped.startswith("### "):
                story.append(Paragraph(stripped[4:], styles["Heading3"]))
                story.append(Spacer(1, 6))
            elif stripped:
                story.append(Paragraph(stripped, styles["Normal"]))
                story.append(Spacer(1, 4))

        doc.build(story)
        logger.info("PDF exported (fallback): %s", out_path)
        return out_path

    # ------------------------------------------------------------------
    # PowerPoint (.pptx)
    # ------------------------------------------------------------------
    def _to_pptx(self, draft: DraftDocument) -> Path:
        out_path = self._outputs_dir / f"{draft.doc_type.value}.pptx"

        # Try OpenDocs
        if self._try_opendocs_pptx(draft, out_path):
            return out_path

        return self._fallback_pptx(draft, out_path)

    def _try_opendocs_pptx(self, draft: DraftDocument, out_path: Path) -> bool:
        try:
            from opendocs.core.parser import ReadmeParser
            from opendocs.generators.pptx_generator import PptxGenerator
            from opendocs.generators.themes import get_theme

            parser = ReadmeParser()
            doc_model = parser.parse(draft.content, repo_name=draft.title)
            theme = get_theme(self._theme_name)
            gen = PptxGenerator(theme=theme)
            result = gen.generate(doc_model, out_path.parent)
            if result.success and result.output_path:
                result_path = Path(result.output_path)
                if result_path != out_path and result_path.exists():
                    result_path.rename(out_path)
                logger.info("PPTX exported via OpenDocs: %s", out_path)
                return True
        except Exception as exc:
            logger.debug("OpenDocs PPTX failed, using fallback: %s", exc)
        return False

    def _fallback_pptx(self, draft: DraftDocument, out_path: Path) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Parse Markdown into slides (split on ##)
        slides_md = self._split_into_slides(draft.content)

        for slide_title, bullets in slides_md:
            layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_title
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(16)

        prs.save(str(out_path))
        logger.info("PPTX exported (fallback): %s", out_path)
        return out_path

    @staticmethod
    def _split_into_slides(md: str) -> list[tuple[str, list[str]]]:
        """Split Markdown into (title, bullets) pairs for slides."""
        slides: list[tuple[str, list[str]]] = []
        current_title = "Untitled"
        bullets: list[str] = []

        for line in md.splitlines():
            stripped = line.strip()
            if stripped.startswith("## ") or stripped.startswith("# "):
                if bullets or current_title != "Untitled":
                    slides.append((current_title, bullets))
                current_title = stripped.lstrip("#").strip()
                bullets = []
            elif stripped.startswith("- ") or stripped.startswith("* "):
                bullets.append(stripped[2:])
            elif stripped.startswith("1. "):
                bullets.append(stripped[3:])
            elif stripped and not stripped.startswith("```"):
                bullets.append(stripped)

        if bullets or current_title != "Untitled":
            slides.append((current_title, bullets))

        return slides or [("Document", ["No content generated"])]
